"""
PowerPoint (.pptx) text extraction.

Extracts text from slides, including titles, body text, notes,
and table content.
"""

import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple


class PPTXExtractor:
    """Extract text and metadata from PowerPoint files."""
    
    def extract(self, file_path: Path) -> List[Dict[str, Any]]:
        """
        Extract text from a PowerPoint file.
        
        Returns one document per presentation (all slides combined).
        
        Args:
            file_path: Path to the .pptx file
            
        Returns:
            List with one dict containing url, title, text, headings, metadata
        """
        try:
            from vendor.pptx import Presentation
        except ImportError:
            try:
                from pptx import Presentation
            except ImportError:
                return [{'error': 'python-pptx not installed. Run: pip install python-pptx'}]
        
        file_path = Path(file_path)
        
        try:
            prs = Presentation(str(file_path))
        except Exception as e:
            return [{'error': f'Failed to open PowerPoint file: {e}'}]
        
        slides_text = []
        headings = []
        total_slides = len(prs.slides)
        
        # Extract metadata
        title = ''
        if prs.core_properties.title:
            title = prs.core_properties.title.strip()
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = []
            slide_title = None
            
            for shape in slide.shapes:
                # Get slide title
                if shape.has_text_frame:
                    if shape.shape_id == 1 or (hasattr(shape, 'placeholder_format') 
                                                and shape.placeholder_format is not None
                                                and shape.placeholder_format.type is not None
                                                and shape.placeholder_format.idx == 0):
                        # Title placeholder
                        text = shape.text_frame.text.strip()
                        if text:
                            slide_title = text
                
                # Extract text from text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            slide_parts.append(' | '.join(row_texts))
            
            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"Notes: {notes_text}")
            
            if slide_title:
                headings.append((2, f"{slide_title}"))
                if not title and slide_num == 1:
                    title = slide_title
            
            if slide_parts:
                slide_text = f"Slide {slide_num}" + (f": {slide_title}" if slide_title else "")
                slide_text += "\n" + "\n".join(slide_parts)
                slides_text.append(slide_text)
        
        full_text = "\n\n".join(slides_text)
        
        # Clean control characters
        full_text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', full_text)
        
        if not title:
            title = file_path.stem
        
        return [{
            'url': file_path.as_uri(),
            'title': title,
            'text': full_text,
            'headings': headings,
            'metadata': {
                'doc_type': 'pptx',
                'total_slides': total_slides,
                'author': prs.core_properties.author or '',
            },
        }]
